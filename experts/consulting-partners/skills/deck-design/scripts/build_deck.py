#!/usr/bin/env python3
"""
⚠️ DEPRECATED(2026-07 废弃,禁止调用)——
本脚本仅有 4 种简陋版式(cover/bullet/chart_bar/two_column),是"产出 PPT 极其简单、大片留白"
事故的根源。已由双引擎 mck_fusion.FusionDeck(mck_ppt 67版式 + mckinsey_pptx 补充版式)取代。
PPT 生成一律走 mck_fusion,见上级目录 SKILL.md。本文件保留仅作历史对照。

PPT 骨架生成器 —— 按统一设计规范批量生成幻灯片框架
原创实现，供 consulting-partners 专家团的交付设计师（deck-designer）使用。

用法：
    python3 build_deck.py --outline outline.json --output deck.pptx

outline.json 格式（骨架确认后由使用者提供）：
[
  {"type": "cover", "title": "XX市场进入战略", "subtitle": "2026年7月"},
  {"type": "bullet", "title": "执行摘要", "bullets": ["要点1", "要点2", "要点3"]},
  {"type": "chart_bar", "title": "市场规模趋势", "categories": ["2023","2024","2025"], "values": [10,15,22], "insight": "洞察文字"},
  {"type": "two_column", "title": "方案对比", "left_title": "方案A", "left_bullets": [...], "right_title": "方案B", "right_bullets": [...]}
]

依赖：pip install python-pptx
"""

import argparse
import json
import sys

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
except ImportError:
    print("请先安装依赖：pip install python-pptx", file=sys.stderr)
    raise

# ---- 设计规范常量（详见 references/design-system.md） ----
NAVY = RGBColor(0x0B, 0x2A, 0x4A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MED_GRAY = RGBColor(0x6E, 0x6E, 0x6E)
ACCENT = RGBColor(0xC8, 0x8A, 0x2E)
BG_LIGHT = RGBColor(0xF4, 0xF4, 0xF2)

TITLE_SIZE = Pt(26)
SUBTITLE_SIZE = Pt(16)
BODY_SIZE = Pt(14)
FOOTNOTE_SIZE = Pt(9)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def _set_title(slide, text, size=TITLE_SIZE, color=NAVY, top=Inches(0.35), left=Inches(0.5)):
    box = slide.shapes.add_textbox(left, top, SLIDE_WIDTH - left * 2, Inches(0.9))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.bold = True
    p.font.color.rgb = color
    return box


def _add_footnote(slide, text="来源：见报告数据来源章节"):
    box = slide.shapes.add_textbox(Inches(0.5), SLIDE_HEIGHT - Inches(0.4), SLIDE_WIDTH - Inches(1), Inches(0.3))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = FOOTNOTE_SIZE
    p.font.color.rgb = MED_GRAY


def build_cover(prs, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(1, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.8), SLIDE_WIDTH - Inches(2), Inches(1.5))
    p = title_box.text_frame.paragraphs[0]
    p.text = page.get("title", "")
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT

    if page.get("subtitle"):
        sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), SLIDE_WIDTH - Inches(2), Inches(0.6))
        sp = sub_box.text_frame.paragraphs[0]
        sp.text = page["subtitle"]
        sp.font.size = SUBTITLE_SIZE
        sp.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    return slide


def build_bullet(prs, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_title(slide, page.get("title", ""))
    box = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), SLIDE_WIDTH - Inches(1.4), Inches(5.3))
    tf = box.text_frame
    tf.word_wrap = True
    bullets = page.get("bullets", [])
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"•  {b}"
        p.font.size = BODY_SIZE
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(10)
    _add_footnote(slide)
    return slide


def build_chart_bar(prs, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_title(slide, page.get("title", ""))

    chart_data = CategoryChartData()
    chart_data.categories = page.get("categories", [])
    chart_data.add_series(page.get("title", "数据"), page.get("values", []))

    x, y, cx, cy = Inches(0.7), Inches(1.5), Inches(7.5), Inches(4.8)
    slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)

    if page.get("insight"):
        insight_box = slide.shapes.add_textbox(Inches(8.5), Inches(1.5), Inches(4.2), Inches(4.8))
        tf = insight_box.text_frame
        tf.word_wrap = True
        bg = slide.shapes.add_shape(1, Inches(8.4), Inches(1.5), Inches(4.4), Inches(4.8))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_LIGHT
        bg.line.color.rgb = ACCENT
        bg.line.width = Pt(1)
        p = tf.paragraphs[0]
        p.text = "洞察"
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = ACCENT
        p2 = tf.add_paragraph()
        p2.text = page["insight"]
        p2.font.size = Pt(12)
        p2.font.color.rgb = DARK_GRAY
    _add_footnote(slide)
    return slide


def build_two_column(prs, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_title(slide, page.get("title", ""))

    col_width = (SLIDE_WIDTH - Inches(1.5)) / 2
    for idx, side in enumerate(["left", "right"]):
        x = Inches(0.5) + idx * (col_width + Inches(0.5))
        header_box = slide.shapes.add_textbox(x, Inches(1.4), col_width, Inches(0.5))
        hp = header_box.text_frame.paragraphs[0]
        hp.text = page.get(f"{side}_title", "")
        hp.font.bold = True
        hp.font.size = Pt(16)
        hp.font.color.rgb = NAVY

        body_box = slide.shapes.add_textbox(x, Inches(2.0), col_width, Inches(4.5))
        tf = body_box.text_frame
        tf.word_wrap = True
        bullets = page.get(f"{side}_bullets", [])
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"•  {b}"
            p.font.size = BODY_SIZE
            p.font.color.rgb = DARK_GRAY
            p.space_after = Pt(8)
    _add_footnote(slide)
    return slide


BUILDERS = {
    "cover": build_cover,
    "bullet": build_bullet,
    "chart_bar": build_chart_bar,
    "two_column": build_two_column,
}


def build_deck(outline: list, output_path: str):
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    for page in outline:
        page_type = page.get("type")
        builder = BUILDERS.get(page_type)
        if not builder:
            print(f"⚠️ 跳过未知版式类型: {page_type}，支持的类型: {list(BUILDERS.keys())}")
            continue
        builder(prs, page)

    prs.save(output_path)
    print(f"✅ 已生成 {len(outline)} 页幻灯片：{output_path}")


def main():
    parser = argparse.ArgumentParser(description="按骨架 JSON 批量生成 PPT 框架")
    parser.add_argument("--outline", required=True, help="骨架 JSON 文件路径")
    parser.add_argument("--output", required=True, help="输出 .pptx 文件路径")
    args = parser.parse_args()

    with open(args.outline, "r", encoding="utf-8") as f:
        outline = json.load(f)

    build_deck(outline, args.output)


if __name__ == "__main__":
    main()
