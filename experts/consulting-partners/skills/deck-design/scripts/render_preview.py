#!/usr/bin/env python3
"""Render a PPTX for visual QC with an explicit, machine-readable mode."""
import argparse
import json
import os
import shutil
import subprocess
import sys
from pathlib import Path

from pptx import Presentation


def _run(command):
    completed = subprocess.run(command, capture_output=True, text=True, check=False)
    if completed.returncode != 0:
        detail = (completed.stderr or completed.stdout or "unknown command failure").strip()
        raise RuntimeError(f"Command failed ({completed.returncode}): {' '.join(command)}: {detail}")
    return completed


def _write_result(path, result):
    Path(path).parent.mkdir(parents=True, exist_ok=True)
    Path(path).write_text(json.dumps(result, ensure_ascii=False, indent=2), encoding="utf-8")


def _write_structure_report(pptx, preview_dir):
    prs = Presentation(str(pptx))
    slides = []
    for slide_num, slide in enumerate(prs.slides, 1):
        texts = []
        chart_count = 0
        table_count = 0
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = (shape.text_frame.text or "").strip()
                if text:
                    texts.append(text)
            chart_count += int(getattr(shape, "has_chart", False))
            table_count += int(getattr(shape, "has_table", False))
        slides.append({
            "slide": slide_num,
            "shape_count": len(slide.shapes),
            "chart_count": chart_count,
            "table_count": table_count,
            "texts": texts,
        })
    report_path = preview_dir / "structure_preview.json"
    _write_result(report_path, {"total_slides": len(slides), "slides": slides})
    return str(report_path)


def render_preview(pptx_path, project_dir):
    pptx = Path(pptx_path).resolve()
    project = Path(project_dir).resolve()
    preview_dir = project / "preview"
    preview_dir.mkdir(parents=True, exist_ok=True)
    result_path = project / "render_result.json"
    result = {
        "passed": False,
        "mode": "unavailable",
        "pptx_path": str(pptx),
        "preview_dir": str(preview_dir),
        "full_render": False,
        "visual_qc_required": True,
        "files": [],
    }

    if not pptx.exists():
        result["error"] = f"PPTX does not exist: {pptx}"
        _write_result(result_path, result)
        return result

    try:
        structure_report = _write_structure_report(pptx, preview_dir)
        result["files"].append(structure_report)
    except Exception as exc:
        result["error"] = f"Cannot inspect PPTX structure: {exc}"
        _write_result(result_path, result)
        return result

    office = shutil.which("libreoffice") or shutil.which("soffice")
    pdftoppm = shutil.which("pdftoppm")
    try:
        if office and pdftoppm:
            _run([office, "--headless", "--convert-to", "pdf", "--outdir", str(preview_dir), str(pptx)])
            pdf_path = preview_dir / f"{pptx.stem}.pdf"
            if not pdf_path.exists():
                raise RuntimeError("LibreOffice did not produce the expected PDF")
            prefix = preview_dir / "slide"
            _run([pdftoppm, "-png", "-r", "144", str(pdf_path), str(prefix)])
            images = sorted(str(path) for path in preview_dir.glob("slide-*.png"))
            if not images:
                raise RuntimeError("PDF renderer did not produce slide images")
            result.update({
                "passed": True,
                "mode": "full",
                "full_render": True,
                "files": result["files"] + [str(pdf_path)] + images,
                "note": "All pages rendered. Human visual QC is still required.",
            })
        elif shutil.which("qlmanage"):
            try:
                _run(["qlmanage", "-t", "-s", "1600", "-o", str(preview_dir), str(pptx)])
                files = sorted(str(path) for path in preview_dir.iterdir() if path.is_file())
                result.update({
                    "passed": True,
                    "mode": "quicklook_first_screen",
                    "full_render": False,
                    "files": files,
                    "note": "Visual gate degraded: Quick Look first screen plus structure report; full-page rendering was not completed.",
                })
            except Exception as exc:
                result.update({
                    "passed": True,
                    "mode": "structure_only",
                    "full_render": False,
                    "renderer_error": str(exc),
                    "note": "Visual gate heavily degraded: structure report only; no rendered preview was available.",
                })
        else:
            result.update({
                "passed": True,
                "mode": "structure_only",
                "full_render": False,
                "note": "Visual gate heavily degraded: structure report only; no rendered preview was available.",
            })
    except Exception as exc:
        result["error"] = str(exc)

    _write_result(result_path, result)
    return result


def main():
    parser = argparse.ArgumentParser(description="Render PPTX previews for visual QC")
    parser.add_argument("pptx", help="Path to the generated PPTX")
    parser.add_argument("project_dir", help="Project directory containing preview/")
    args = parser.parse_args()
    result = render_preview(args.pptx, args.project_dir)
    print(f"[render_preview] mode={result['mode']} full_render={result['full_render']}")
    if result.get("note"):
        print(f"[render_preview] {result['note']}")
    if result.get("error"):
        print(f"[render_preview] {result['error']}", file=sys.stderr)
    sys.exit(0 if result["passed"] else 1)


if __name__ == "__main__":
    main()
